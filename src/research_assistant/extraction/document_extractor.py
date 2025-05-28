import docx
from typing import Dict, Any, Optional, List
from datetime import datetime
import os
from pathlib import Path
import mimetypes
import asyncio
import re
import langdetect
from pptx import Presentation
import mammoth
import python_pptx

from research_assistant.extraction.base_extractor import BaseExtractor, ExtractedContent, ExtractionOptions
from research_assistant.utils.logging import get_logger

logger = get_logger(__name__)

class DocumentExtractor(BaseExtractor):
    """Document content extractor implementation for Word and PowerPoint files."""

    def __init__(self):
        """Initialize the document extractor."""
        super().__init__(
            name="document",
            description="Extract content from Word and PowerPoint files"
        )

    async def extract(
        self,
        source: str,
        options: Optional[ExtractionOptions] = None
    ) -> ExtractedContent:
        """
        Extract content from a document file.

        Args:
            source: Path to the document file
            options: Optional extraction options

        Returns:
            Extracted content

        Raises:
            ValueError: If file is invalid
            Exception: For other extraction errors
        """
        if not await self.validate_source(source):
            raise ValueError(f"Invalid document file: {source}")

        options = options or ExtractionOptions()
        path = Path(source)
        suffix = path.suffix.lower()

        try:
            # Check file size
            if options.max_size and path.stat().st_size > options.max_size:
                raise ValueError(f"File size exceeds maximum limit: {path.stat().st_size} bytes")

            # Extract content based on file type
            if suffix == '.docx':
                content, metadata = await self._extract_docx(path, options)
            elif suffix == '.pptx':
                content, metadata = await self._extract_pptx(path, options)
            else:
                raise ValueError(f"Unsupported file type: {suffix}")

            # Clean content
            content = self._clean_content(content)

            # Detect language
            try:
                language = langdetect.detect(content)
            except:
                language = None

            # Count words and characters
            word_count = len(content.split())
            char_count = len(content)

            return self.format_content({
                "title": metadata.get("title", path.stem),
                "text": content,
                "metadata": metadata,
                "language": language,
                "word_count": word_count,
                "char_count": char_count
            })

        except Exception as e:
            self.logger.error(f"Error extracting from {source}: {str(e)}")
            raise

    async def validate_source(self, source: str) -> bool:
        """
        Validate if the source is a valid document file.

        Args:
            source: File path to validate

        Returns:
            True if file is valid, False otherwise
        """
        try:
            path = Path(source)
            return (
                path.exists() and 
                path.is_file() and 
                path.suffix.lower() in ['.docx', '.pptx']
            )
        except:
            return False

    async def _extract_docx(
        self,
        path: Path,
        options: ExtractionOptions
    ) -> tuple[str, Dict[str, Any]]:
        """
        Extract content from a Word document.

        Args:
            path: Path to the Word document
            options: Extraction options

        Returns:
            Tuple of (content, metadata)
        """
        try:
            # Load document
            doc = docx.Document(path)
            
            # Extract metadata
            metadata = {
                "title": "",
                "author": "",
                "created": None,
                "modified": None,
                "pages": 0,
                "sections": len(doc.sections),
                "tables": len(doc.tables),
                "images": 0
            }

            # Extract core properties
            core_props = doc.core_properties
            if core_props:
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "created": core_props.created,
                    "modified": core_props.modified
                })

            # Extract content
            content_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)

            # Extract tables if requested
            if options.extract_tables:
                for table in doc.tables:
                    table_text = []
                    for row in table.rows:
                        row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        content_parts.append("\n".join(table_text))

            # Count images
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    metadata["images"] += 1

            return "\n\n".join(content_parts), metadata

        except Exception as e:
            self.logger.error(f"Error extracting from Word document: {str(e)}")
            raise

    async def _extract_pptx(
        self,
        path: Path,
        options: ExtractionOptions
    ) -> tuple[str, Dict[str, Any]]:
        """
        Extract content from a PowerPoint presentation.

        Args:
            path: Path to the PowerPoint file
            options: Extraction options

        Returns:
            Tuple of (content, metadata)
        """
        try:
            # Load presentation
            prs = Presentation(path)
            
            # Extract metadata
            metadata = {
                "title": "",
                "author": "",
                "created": None,
                "modified": None,
                "slides": len(prs.slides),
                "images": 0
            }

            # Extract core properties
            core_props = prs.core_properties
            if core_props:
                metadata.update({
                    "title": core_props.title or "",
                    "author": core_props.author or "",
                    "created": core_props.created,
                    "modified": core_props.modified
                })

            # Extract content
            content_parts = []

            # Extract slides
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"Slide {slide_num}:"]

                # Extract shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text)

                    # Count images
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        metadata["images"] += 1

                if len(slide_parts) > 1:
                    content_parts.append("\n".join(slide_parts))

            # Extract notes if available
            if options.include_metadata:
                for slide_num, slide in enumerate(prs.slides, 1):
                    if slide.has_notes_slide:
                        notes = slide.notes_slide.notes_text_frame.text
                        if notes.strip():
                            content_parts.append(f"Notes for Slide {slide_num}:\n{notes}")

            return "\n\n".join(content_parts), metadata

        except Exception as e:
            self.logger.error(f"Error extracting from PowerPoint: {str(e)}")
            raise

    def _clean_content(self, content: str) -> str:
        """
        Clean extracted content.

        Args:
            content: Raw content to clean

        Returns:
            Cleaned content
        """
        # Remove extra whitespace
        content = re.sub(r'\s+', ' ', content)
        
        # Remove empty lines
        content = re.sub(r'\n\s*\n', '\n', content)
        
        # Remove special characters
        content = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]', '', content)
        
        return content.strip()

    async def close(self) -> None:
        """Close the document extractor."""
        pass  # No resources to clean up 